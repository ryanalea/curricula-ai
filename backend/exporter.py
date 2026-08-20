import json
import io
import zipfile
import re
import html as html_lib

try:
    import docx
    from docx.shared import Inches, Pt, RGBColor
except ImportError:
    docx = None

try:
    import pdfkit
except ImportError:
    pdfkit = None

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    sync_playwright = None

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
except ImportError:
    SimpleDocTemplate = None

def get_role_label(role: str) -> str:
    mapping = {
        "creator": "Creator POV",
        "student": "Student POV",
        "educator": "Educator POV",
        "all": "All Roles Combined"
    }
    return mapping.get(role.lower(), role.capitalize())

def clean_lesson_title(title: str) -> str:
    if not title:
        return "Untitled Lesson"
    cleaned = re.sub(r'^\s*Lesson\s*\d+\s*[:\-\.]*\s*', '', title, flags=re.IGNORECASE).strip()
    return cleaned if cleaned else title

def format_section_content_to_md(content, indent: int = 0) -> list[str]:
    lines = []
    pad = "  " * indent

    if isinstance(content, str):
        lines.append(f"{pad}{content}")
    elif isinstance(content, list):
        for item in content:
            if isinstance(item, dict):
                if "question" in item:  # Quiz item
                    q = item.get("question", "")
                    lines.append(f"{pad}- **Question:** {q}")
                    if "options" in item and isinstance(item["options"], list):
                        for opt in item["options"]:
                            correct_marker = " **(Correct Answer)**" if opt == item.get("answer") else ""
                            lines.append(f"{pad}  - {opt}{correct_marker}")
                    if item.get("explanation"):
                        lines.append(f"{pad}  - *Explanation:* {item['explanation']}")
                elif "criteria" in item:  # Rubric item
                    c = item.get("criteria", "")
                    lines.append(f"{pad}- **Criteria:** {c}")
                    for k in ["excellent", "good", "needs_improvement"]:
                        if item.get(k):
                            lines.append(f"{pad}  - *{k.replace('_', ' ').capitalize()}:* {item[k]}")
                elif "title" in item or "name" in item:  # Exercise or activity item
                    t = item.get("title") or item.get("name")
                    lines.append(f"{pad}- **{t}**")
                    for k, v in item.items():
                        if k in ["title", "name"]:
                            continue
                        if isinstance(v, (str, int, float)):
                            lines.append(f"{pad}  - *{k.replace('_', ' ').capitalize()}:* {v}")
                        elif isinstance(v, list):
                            lines.append(f"{pad}  - *{k.replace('_', ' ').capitalize()}:* {', '.join(map(str, v))}")
                else:
                    parts = []
                    for k, v in item.items():
                        if isinstance(v, (str, int, float)):
                            parts.append(f"*{k.replace('_', ' ').capitalize()}:* {v}")
                    if parts:
                        lines.append(f"{pad}- " + " | ".join(parts))
                    else:
                        lines.append(f"{pad}- {json.dumps(item)}")
            else:
                lines.append(f"{pad}- {item}")
    elif isinstance(content, dict):
        for k, v in content.items():
            k_label = k.replace('_', ' ').capitalize()
            if isinstance(v, (str, int, float)):
                lines.append(f"{pad}**{k_label}:** {v}")
            elif isinstance(v, (list, dict)):
                lines.append(f"{pad}**{k_label}:**")
                lines.extend(format_section_content_to_md(v, indent + 1))
            else:
                lines.append(f"{pad}**{k_label}:** {v}")
    else:
        lines.append(f"{pad}{content}")
    return lines

def get_resolved_lesson_sections(lesson: dict, role: str) -> dict:
    sections = lesson.get("sections", {}).get(role, {})
    if sections and len(sections) > 0:
        return sections
        
    title = clean_lesson_title(lesson.get("title", "Lesson Content"))
    if role == "creator":
        return {
            "overview": f"This lesson provides a comprehensive overview and practical foundation for {title}. Students will explore core concepts, industry use-cases, and implementation patterns necessary for real-world projects.",
            "learning_outcomes": [
                f"Master core concepts and architectural components of {title}.",
                f"Implement hands-on code examples and workflows using industry standards.",
                "Apply critical thinking to analyze, debug, and optimize real-world production scenarios."
            ],
            "core_content": f"### 1. Conceptual Foundations\n{title} serves as a key pillar in modern systems engineering. By leveraging structured workflows and robust error handling, developers can ensure high performance and maintainability.\n\n### 2. Practical Implementation\nTo implement {title} effectively, engineers must follow clean architecture patterns and best practices.",
            "exercises": [
                {"title": f"Building {title} Pipeline", "description": f"Implement a basic working prototype for {title} using Python/JavaScript.", "code_template": f"// Exercise: {title}\nfunction executeTask() {{\n  console.log('Executing {title}...');\n}}"}
            ],
            "quizzes": [
                {"question": f"What is the primary objective of {title}?", "options": ["To establish a robust, scalable technical workflow", "To bypass data validation", "To reduce readability"], "answer": "To establish a robust, scalable technical workflow", "explanation": "It ensures reliable engineering standards."}
            ]
        }
    elif role == "student":
        return {
            "why_this_matters": f"Understanding {title} is crucial for career advancement. It bridges theoretical principles with industry-grade implementation strategies.",
            "practice": {
                "code_block": f"// Interactive Sandbox for {title}\nfunction main() {{\n  console.log('Running {title} sandbox...');\n}}\nmain();",
                "interactive_exercise": f"Extend the function logic for {title}.",
                "checklist": ["Set up local environment", "Implement core logic", "Pass automated tests"]
            },
            "debugging": "### Common Pitfalls & Solutions\n1. **Unhandled Edge Cases:** Validate inputs prior to execution.\n2. **Performance Bottlenecks:** Optimize data structure lookups.",
            "ethics": "### Code Principles & Ethics\nEnsure user data protection, transparency, and security compliance throughout implementation."
        }
    else: # educator
        return {
            "facilitator_guide": f"### Educator Instructions\nFacilitate an interactive discussion on {title}. Encourage students to participate in pair-programming exercises.",
            "lesson_plan": {
                "ice_breaker": f"Ask students: 'What real-world applications of {title} have you encountered?'",
                "timing": "Lecture & Demo: 20 mins | Pair Lab: 30 mins | Wrap-up & Q&A: 10 mins"
            },
            "rubric": [
                {"criteria": "Implementation", "excellent": "Code runs error-free with optimal logic", "good": "Code runs with minor style issues", "needs_improvement": "Code contains execution errors"},
                {"criteria": "Understanding", "excellent": "Demonstrates deep mastery of concepts", "good": "Demonstrates basic understanding", "needs_improvement": "Lacks core understanding"}
            ],
            "discussion_questions": [
                f"How does {title} improve overall system efficiency?",
                "What key trade-offs should be considered when deploying to production?"
            ]
        }

def export_to_markdown(course_data: dict, role: str) -> str:
    md = []
    md.append(f"# 🎓 {course_data.get('title', 'Untitled Course')}")
    md.append(f"**Difficulty:** {course_data.get('config', {}).get('difficulty', 'Beginner')} | **Audience:** {course_data.get('config', {}).get('target_audience', 'Student')}\n")
    md.append("---\n")

    roles_to_export = ["creator", "student", "educator"] if role == "all" else [role]
    lessons = course_data.get("lessons", [])
    if not lessons and course_data.get("structure"):
        lessons = course_data.get("structure")

    if not lessons:
        lessons = [{"title": course_data.get("title", "Course Module")}]

    for r in roles_to_export:
        md.append(f"## 📘 {get_role_label(r)}")
        for idx, lesson in enumerate(lessons):
            l_num = lesson.get('order') or (idx + 1)
            clean_t = clean_lesson_title(lesson.get('title', 'Untitled Lesson'))
            md.append(f"### Lesson {l_num}: {clean_t}")
            sections = get_resolved_lesson_sections(lesson, r)
            
            for sec_type, content in sections.items():
                title = sec_type.replace("_", " ").capitalize()
                md.append(f"#### {title}")
                formatted_lines = format_section_content_to_md(content)
                md.extend(formatted_lines)
                md.append("")
        md.append("---\n")
    return "\n".join(md)

def _inline_markdown_to_html(text: str) -> str:
    """Converts inline markdown (bold, italics, inline code) to HTML."""
    # Inline code first, so markers inside `code` aren't picked up as bold/italic
    text = re.sub(r'`([^`]+?)`', r"<code style='background:#F1F5F9;color:#BE185D;padding:1px 5px;border-radius:4px;font-family:Consolas,monospace;font-size:0.9em;'>\1</code>", text)
    # Bold (**text** or __text__)
    text = re.sub(r'\*\*(.+?)\*\*', r"<strong>\1</strong>", text)
    text = re.sub(r'__(.+?)__', r"<strong>\1</strong>", text)
    # Italics (*text* or _text_) - avoid matching leftover single asterisks used as bullets
    text = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r"<em>\1</em>", text)
    return text

def _md_block_to_html(md_text: str) -> str:
    """Renders a markdown snippet (paragraphs, ### headers, - lists, ```code```)
    into the dark-theme component markup used by the new PDF/HTML template."""
    if not md_text:
        return ""
    md_text = re.sub(r'[\U00010000-\U0010ffff]', '', str(md_text))
    lines_out = []
    in_code = False
    code_buf = []
    list_open = False

    def close_list():
        nonlocal list_open
        if list_open:
            lines_out.append("</ul>")
            list_open = False

    for raw in str(md_text).split("\n"):
        line = raw.strip()
        if line.startswith("```"):
            if not in_code:
                close_list()
                in_code = True
                code_buf = []
            else:
                in_code = False
                code_text = html_lib.escape("\n".join(code_buf))
                lines_out.append(f"<div class='codeblock'>{code_text}</div>")
            continue
        if in_code:
            code_buf.append(raw)
            continue
        if line.startswith("### "):
            close_list()
            lines_out.append(f"<div class='md-h3'>{_inline_markdown_to_html(html_lib.escape(line[4:]))}</div>")
        elif line.startswith("#### "):
            close_list()
            lines_out.append(f"<div class='md-h4'>{_inline_markdown_to_html(html_lib.escape(line[5:]))}</div>")
        elif line.startswith("- ") or line.startswith("* "):
            if not list_open:
                lines_out.append("<ul class='md-ul'>")
                list_open = True
            lines_out.append(f"<li>{_inline_markdown_to_html(html_lib.escape(line[2:]))}</li>")
        elif line == "":
            close_list()
        else:
            close_list()
            lines_out.append(f"<p class='md-p'>{_inline_markdown_to_html(html_lib.escape(line))}</p>")
    close_list()
    return "\n".join(lines_out)


# ---------------------------------------------------------------------------
# New dark-theme, component-based HTML/PDF template
# (cover page, TOC, quiz cards, exercise cards, code blocks, comparison-style
# cards) rendered through Chromium (Playwright) for full modern CSS support.
# ---------------------------------------------------------------------------

_TEMPLATE_CSS = """
@page { size: A4; margin: 0; }
* { box-sizing: border-box; }
body {
    margin: 0; background: #ffffff; color: #1e2433;
    font-family: 'Helvetica Neue', Arial, sans-serif;
}
.cover, .toc-page {
    width: 210mm; min-height: 297mm; display: flex; flex-direction: column; justify-content: space-between;
    page-break-after: always; background: #ffffff;
    padding: 18mm 16mm;
}
.content-wrap { padding: 16mm 16mm 4mm; }
.brand { display: flex; align-items: center; gap: 12px; }
.brand-badge {
    width: 38px; height: 38px; border-radius: 10px;
    background: linear-gradient(135deg,#2563eb,#0d9488);
    display: flex; align-items: center; justify-content: center;
    font-weight: 800; color: #ffffff; font-size: 15px;
}
.brand-name { font-size: 15px; font-weight: 800; color: #1e2433; line-height: 1.15; }
.brand-sub { font-size: 9px; letter-spacing: 2px; color: #0d9488; font-weight: 700; }
.cover-top { display: flex; justify-content: space-between; align-items: flex-start; }
.module-badge {
    border: 1px solid rgba(15,23,42,0.12); border-radius: 8px;
    padding: 8px 14px; font-size: 9px; letter-spacing: 2px; color: #5b6178;
    text-align: right; text-transform: uppercase; line-height: 1.6;
}
.eyebrow { font-size: 10px; letter-spacing: 3px; text-transform: uppercase; color: #0d9488; font-weight: 700; }
.cover h1 { font-size: 40px; line-height: 1.15; font-weight: 800; color: #1e2433; margin: 14px 0 0; max-width: 500px; }
.cover h1 .accent { color: #2563eb; }
.cover-desc { margin-top: 18px; font-size: 13px; line-height: 1.7; color: #5b6178; max-width: 460px; }
.cover-desc b { color: #1e2433; }
.pill-row { margin-top: 22px; display: flex; flex-wrap: wrap; gap: 8px; }
.pill {
    display: inline-flex; align-items: center; gap: 6px;
    border: 1px solid rgba(15,23,42,0.12); background: #f8fafc;
    border-radius: 999px; padding: 7px 14px; font-size: 10.5px; color: #3b4257;
}
.pill .dot { width: 6px; height: 6px; border-radius: 50%; background: #16a34a; }
.pill .dot.blue { background: #2563eb; }
.cover-footer, .toc-footer {
    display: flex; justify-content: space-between; border-top: 1px solid rgba(15,23,42,0.1); padding-top: 14px;
}
.cover-footer .label, .toc-footer .label { font-size: 9px; letter-spacing: 1px; color: #8b93b8; }
.cover-footer .value, .toc-footer .value { font-size: 12px; color: #1e2433; font-weight: 600; margin-top: 4px; }

.toc-title { font-size: 28px; font-weight: 800; color: #1e2433; margin: 10px 0 22px; }
.stat-row { display: grid; grid-template-columns: repeat(4,1fr); gap: 10px; margin-bottom: 26px; }
.stat-card { border: 1px solid rgba(15,23,42,0.1); border-radius: 10px; padding: 12px; background: #f8fafc; }
.stat-card .num { font-size: 13px; font-weight: 800; color: #2563eb; }
.stat-card .lbl { font-size: 8.5px; letter-spacing: 1px; color: #7a81a3; margin-top: 4px; text-transform: uppercase; }
.toc-row {
    display: flex; align-items: center; gap: 14px; padding: 11px 4px;
    border-bottom: 1px solid rgba(15,23,42,0.08); font-size: 12.5px; color: #2b3147;
}
.toc-num { color: #8b93b8; font-weight: 700; width: 20px; }
.toc-row .flex1 { flex: 1; }
.toc-row .pgnum { color: #0d9488; font-weight: 700; }

/* ---- flowing content pages ---- */
.content-wrap { background: #ffffff; }
.lesson-divider {
    page-break-before: always; padding-top: 4mm;
}
.kicker-row { display: flex; align-items: center; gap: 10px; margin-bottom: 14px; }
.kicker-num {
    width: 24px; height: 24px; border-radius: 6px; background: rgba(13,148,136,0.12);
    color: #0d9488; font-size: 11px; font-weight: 800; display: flex; align-items: center; justify-content: center;
}
.section-tag { font-size: 10px; letter-spacing: 3px; text-transform: uppercase; color: #7a81a3; font-weight: 700; }
.role-tag {
    display: inline-block; margin-bottom: 8px; font-size: 9px; letter-spacing: 2px; text-transform: uppercase;
    color: #ffffff; background: #0d9488; border-radius: 4px; padding: 3px 9px; font-weight: 800;
}
h2.section-title { font-size: 22px; font-weight: 800; color: #1e2433; margin: 4px 0 16px; page-break-after: avoid; }
h3.lesson-title { font-size: 26px; font-weight: 800; color: #1e2433; margin: 6px 0 4px; }

.md-h3 { font-size: 14px; font-weight: 700; color: #1e2433; margin: 16px 0 6px; page-break-after: avoid; }
.md-h4 { font-size: 12px; font-weight: 700; color: #1d4ed8; margin: 12px 0 4px; text-transform: uppercase; letter-spacing: 0.5px; }
.md-p { font-size: 11.5px; color: #3b4257; line-height: 1.7; margin: 0 0 8px; }
.md-ul { margin: 4px 0 10px; padding-left: 18px; }
.md-ul li { font-size: 11.5px; color: #3b4257; line-height: 1.7; margin-bottom: 3px; }

.codeblock {
    background: #f8fafc; border: 1px solid rgba(15,23,42,0.1); border-radius: 10px;
    padding: 12px 14px; font-family: Consolas, 'Liberation Mono', Menlo, monospace;
    font-size: 9.5px; line-height: 1.6; color: #1e2433; white-space: pre-wrap; word-wrap: break-word;
    margin: 8px 0 12px; page-break-inside: avoid;
}

.check-list { display: flex; flex-direction: column; gap: 8px; margin: 6px 0 14px; }
.check-item { display: flex; gap: 10px; align-items: flex-start; font-size: 11.5px; color: #2b3147; line-height: 1.6; page-break-inside: avoid; }
.check-item .tick {
    flex: none; width: 17px; height: 17px; border-radius: 50%; background: rgba(22,163,74,0.12);
    color: #16a34a; font-size: 10px; font-weight: 800; display: flex; align-items: center; justify-content: center; margin-top: 1px;
}

.card { border: 1px solid rgba(15,23,42,0.1); border-radius: 12px; padding: 16px 18px; margin-bottom: 12px; background: #f8fafc; page-break-inside: avoid; }
.ex-head { display: flex; justify-content: space-between; align-items: center; margin-bottom: 8px; }
.ex-head .t { font-size: 13px; font-weight: 700; color: #1e2433; }
.badge { font-size: 8.5px; font-weight: 800; letter-spacing: 1px; padding: 3px 9px; border-radius: 999px; text-transform: uppercase; }
.badge.medium { background: rgba(202,138,4,0.15); color: #a16207; }
.badge.hard { background: rgba(220,38,38,0.15); color: #b91c1c; }
.badge.easy { background: rgba(22,163,74,0.15); color: #15803d; }
.card-desc { font-size: 11.5px; color: #5b6178; line-height: 1.65; }

.quiz-q { font-size: 12.5px; font-weight: 700; color: #1e2433; margin-bottom: 10px; }
.quiz-opt { display: flex; align-items: center; gap: 9px; font-size: 11px; color: #3b4257; padding: 4px 0; }
.quiz-opt .radio { width: 12px; height: 12px; border-radius: 50%; border: 1.5px solid #aab1ce; flex: none; }
.quiz-opt.correct .radio { background: #16a34a; border-color: #16a34a; }
.quiz-opt.correct { color: #1e2433; font-weight: 600; }
.quiz-why { margin-top: 8px; font-size: 10.5px; color: #2563eb; line-height: 1.6; }

table.cmp { width: 100%; border-collapse: collapse; font-size: 10.5px; margin: 8px 0 14px; page-break-inside: avoid; }
table.cmp th { text-align: left; color: #0d9488; font-size: 9px; letter-spacing: 1px; text-transform: uppercase; padding: 8px 10px; border-bottom: 1px solid rgba(15,23,42,0.15); }
table.cmp td { padding: 9px 10px; border-bottom: 1px solid rgba(15,23,42,0.08); color: #3b4257; line-height: 1.5; vertical-align: top; }
table.cmp tr td:first-child { color: #2563eb; font-weight: 700; }

.doc-footer { text-align: center; font-size: 9px; color: #8b93b8; margin-top: 20px; padding-top: 10px; border-top: 1px solid rgba(15,23,42,0.08); }
"""


def _pill(label: str, blue: bool = False) -> str:
    dot_cls = "dot blue" if blue else "dot"
    return f"<span class='pill'><span class='{dot_cls}'></span>{html_lib.escape(str(label))}</span>"


def _render_quiz_card(item: dict) -> str:
    q = html_lib.escape(str(item.get("question", "")))
    options = item.get("options", []) or []
    answer = item.get("answer")
    opts_html = []
    for opt in options:
        is_correct = (opt == answer)
        cls = "quiz-opt correct" if is_correct else "quiz-opt"
        opts_html.append(f"<div class='{cls}'><div class='radio'></div>{html_lib.escape(str(opt))}</div>")
    why = item.get("explanation", "")
    why_html = f"<div class='quiz-why'>Why: {html_lib.escape(str(why))}</div>" if why else ""
    return (
        "<div class='card'>"
        f"<div class='quiz-q'>{q}</div>"
        + "".join(opts_html)
        + why_html +
        "</div>"
    )


def _render_exercise_card(item: dict) -> str:
    title = html_lib.escape(str(item.get("title") or item.get("name") or "Exercise"))
    desc = html_lib.escape(str(item.get("description", "")))
    difficulty = str(item.get("difficulty", "")).lower()
    badge_html = f"<span class='badge {difficulty}'>{html_lib.escape(difficulty)}</span>" if difficulty in ("easy", "medium", "hard") else ""
    code = item.get("code_template") or item.get("starter_code")
    code_html = f"<div class='codeblock'>{html_lib.escape(str(code))}</div>" if code else ""
    return (
        "<div class='card'>"
        f"<div class='ex-head'><div class='t'>{title}</div>{badge_html}</div>"
        f"<div class='card-desc'>{desc}</div>"
        f"{code_html}"
        "</div>"
    )


def _render_rubric_table(rubric: list) -> str:
    rows = []
    for r in rubric:
        crit = html_lib.escape(str(r.get("criteria", "")))
        exc = html_lib.escape(str(r.get("excellent", "")))
        good = html_lib.escape(str(r.get("good", "")))
        needs = html_lib.escape(str(r.get("needs_improvement", "")))
        rows.append(f"<tr><td>{crit}</td><td>{exc}</td><td>{good}</td><td>{needs}</td></tr>")
    return (
        "<table class='cmp'><tr><th>Criteria</th><th>Excellent</th><th>Good</th><th>Needs Improvement</th></tr>"
        + "".join(rows) + "</table>"
    )


def _render_section(section_type: str, content) -> str:
    """Renders one lesson section (by its known type) into dark-theme HTML."""
    label = section_type.replace("_", " ").title()
    out = [f"<div class='md-h3' style='font-size:15px;margin-top:20px;'>{html_lib.escape(label)}</div>"]

    if section_type == "learning_outcomes" and isinstance(content, list):
        items = "".join(
            f"<div class='check-item'><div class='tick'>&#10003;</div>{html_lib.escape(str(c))}</div>"
            for c in content
        )
        out.append(f"<div class='check-list'>{items}</div>")
    elif section_type == "exercises" and isinstance(content, list):
        out.extend(_render_exercise_card(item) if isinstance(item, dict) else f"<div class='card'><div class='card-desc'>{html_lib.escape(str(item))}</div></div>" for item in content)
    elif section_type == "quizzes" and isinstance(content, list):
        out.extend(_render_quiz_card(item) if isinstance(item, dict) else "" for item in content)
    elif section_type == "rubric" and isinstance(content, list):
        out.append(_render_rubric_table(content))
    elif section_type == "discussion_questions" and isinstance(content, list):
        items = "".join(f"<li>{html_lib.escape(str(c))}</li>" for c in content)
        out.append(f"<ul class='md-ul'>{items}</ul>")
    elif section_type == "practice" and isinstance(content, dict):
        code = content.get("code_block")
        if code:
            out.append(f"<div class='codeblock'>{html_lib.escape(str(code))}</div>")
        if content.get("interactive_exercise"):
            out.append(f"<p class='md-p'><b>Try it:</b> {html_lib.escape(str(content['interactive_exercise']))}</p>")
        checklist = content.get("checklist")
        if isinstance(checklist, list) and checklist:
            items = "".join(
                f"<div class='check-item'><div class='tick'>&#10003;</div>{html_lib.escape(str(c))}</div>"
                for c in checklist
            )
            out.append(f"<div class='check-list'>{items}</div>")
    elif section_type == "lesson_plan" and isinstance(content, dict):
        for k, v in content.items():
            out.append(f"<p class='md-p'><b>{html_lib.escape(k.replace('_',' ').title())}:</b> {html_lib.escape(str(v))}</p>")
    elif isinstance(content, str):
        out.append(_md_block_to_html(content))
    elif isinstance(content, list):
        items = "".join(f"<li>{html_lib.escape(str(c)) if not isinstance(c, dict) else html_lib.escape(json.dumps(c))}</li>" for c in content)
        out.append(f"<ul class='md-ul'>{items}</ul>")
    elif isinstance(content, dict):
        for k, v in content.items():
            if isinstance(v, (str, int, float)):
                out.append(f"<p class='md-p'><b>{html_lib.escape(k.replace('_',' ').title())}:</b> {html_lib.escape(str(v))}</p>")
    return "\n".join(out)


def export_to_html_v2(course_data: dict, role: str) -> str:
    """New dark-theme, component-based HTML template matching the
    Maxy Academy / Curricula AI PDF design (cover, TOC, quiz cards,
    exercise cards, code blocks)."""
    title = course_data.get("title", "Untitled Course")
    cfg = course_data.get("config", {})
    difficulty = cfg.get("difficulty", "Beginner")
    audience = cfg.get("target_audience", "Student")
    lessons = course_data.get("lessons", []) or course_data.get("structure") or [{"title": title}]
    roles_to_export = ["creator", "student", "educator"] if role == "all" else [role]

    parts = []

    # ---- Cover ----
    parts.append(f"""
    <div class="cover">
        <div class="cover-top">
            <div class="brand">
                <div class="brand-badge">M</div>
                <div><div class="brand-name">Maxy<br>Academy</div><div class="brand-sub">CURRICULA AI</div></div>
            </div>
            <div class="module-badge">{html_lib.escape(get_role_label(role).upper())}<br>{len(lessons)} LESSON(S)</div>
        </div>
        <div>
            <div class="eyebrow">&mdash; &nbsp;COURSE MATERIAL</div>
            <h1>{_inline_markdown_to_html(html_lib.escape(title))}</h1>
            <div class="cover-desc">Generated course material covering <b>{len(lessons)} lesson(s)</b> for {html_lib.escape(str(audience))} at {html_lib.escape(str(difficulty))} level.</div>
            <div class="pill-row">
                {_pill('Difficulty: ' + str(difficulty))}
                {_pill('Audience: ' + str(audience), blue=True)}
                {_pill('Format: ' + get_role_label(role))}
            </div>
        </div>
        <div class="cover-footer">
            <div><div class="label">AUDIENCE</div><div class="value">{html_lib.escape(str(audience))}</div></div>
            <div style="text-align:right;"><div class="label">PREPARED FOR</div><div class="value">Curricula AI Creator Program</div></div>
        </div>
    </div>
    """)

    # ---- TOC ----
    toc_rows = "".join(
        f"<div class='toc-row'><div class='toc-num'>{(l.get('order') or i+1):02d}</div><div class='flex1'>{html_lib.escape(clean_lesson_title(l.get('title','Untitled Lesson')))}</div></div>"
        for i, l in enumerate(lessons)
    )
    parts.append(f"""
    <div class="toc-page">
        <div>
            <div class="eyebrow">IN THIS DOCUMENT</div>
            <div class="toc-title">Table of contents</div>
            <div class="stat-row">
                <div class="stat-card"><div class="num">{html_lib.escape(str(difficulty))}</div><div class="lbl">Difficulty</div></div>
                <div class="stat-card"><div class="num">{len(lessons)}</div><div class="lbl">Lessons</div></div>
                <div class="stat-card"><div class="num">{html_lib.escape(str(audience))}</div><div class="lbl">Audience</div></div>
                <div class="stat-card"><div class="num">{html_lib.escape(get_role_label(role))}</div><div class="lbl">Perspective</div></div>
            </div>
            {toc_rows}
        </div>
        <div class="toc-footer">
            <div><div class="label">DOCUMENT</div><div class="value">{html_lib.escape(title)}</div></div>
            <div style="text-align:right;"><div class="label">SOURCE</div><div class="value">Curricula AI</div></div>
        </div>
    </div>
    """)

    # ---- Lesson content ----
    content_html = ["<div class='content-wrap'>"]
    for idx, lesson in enumerate(lessons):
        lesson_title = clean_lesson_title(lesson.get("title", "Untitled Lesson"))
        lesson_num = lesson.get('order') or (idx + 1)
        content_html.append(f"""
        <div class="lesson-divider">
            <div class="kicker-row"><div class="kicker-num">{lesson_num:02d}</div><div class="section-tag">LESSON {lesson_num:02d}</div></div>
            <h3 class="lesson-title">{_inline_markdown_to_html(html_lib.escape(lesson_title))}</h3>
        </div>
        """)
        for r in roles_to_export:
            sections = get_resolved_lesson_sections(lesson, r)
            if not sections:
                continue
            if len(roles_to_export) > 1:
                content_html.append(f"<div class='role-tag'>{html_lib.escape(get_role_label(r))}</div>")
            for sec_type, content in sections.items():
                content_html.append(_render_section(sec_type, content))

    content_html.append("</div>")
    content_html.append(f"<div class='doc-footer'>Maxy Academy &middot; Curricula AI &middot; {html_lib.escape(get_role_label(role))}</div>")
    parts.append("\n".join(content_html))

    body = "\n".join(parts)
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>{html_lib.escape(title)}</title>
<style>{_TEMPLATE_CSS}</style>
</head>
<body>
{body}
</body>
</html>"""


def export_to_html(course_data: dict, role: str) -> str:
    """Public entry point used by main.py — renders the new dark-theme
    template. Falls back to the legacy markdown-based renderer only if
    something in the structured renderer blows up on malformed data."""
    try:
        return export_to_html_v2(course_data, role)
    except Exception as e:
        print(f"New HTML template failed, falling back to legacy renderer: {e}")
        return export_to_html_legacy(course_data, role)


def export_to_html_legacy(course_data: dict, role: str) -> str:
    md_content = export_to_markdown(course_data, role)

    # Sanitize emojis that wkhtmltopdf cannot render in default Windows fonts
    md_content = re.sub(r'[\U00010000-\U0010ffff]', '', md_content)
    md_content = re.sub(r'[\u2600-\u27BF]', '', md_content)

    _pdf_difficulty = course_data.get('config', {}).get('difficulty', 'Beginner')
    _pdf_audience = course_data.get('config', {}).get('target_audience', 'Student')

    html_lines = []
    in_code_block = False
    code_buffer = []
    list_mode = None  # None | "ul" | "ul-nested"

    def close_list():
        nonlocal list_mode
        if list_mode == "ul":
            html_lines.append("</ul>")
        elif list_mode == "ul-nested":
            html_lines.append("</ul></li></ul>")
        list_mode = None

    lines = md_content.split("\n")
    for raw_line in lines:
        line_clean = raw_line.strip()

        # --- Fenced code blocks (```lang ... ```) ---
        if line_clean.startswith("```"):
            if not in_code_block:
                close_list()
                in_code_block = True
                code_buffer = []
            else:
                in_code_block = False
                code_text = html_lib.escape("\n".join(code_buffer))
                html_lines.append(
                    "<div style='margin:12px 0;border-radius:8px;overflow:hidden;border:1px solid #334155;'>"
                    "<div style='background:#0F172A;padding:6px 12px;display:block;'>"
                    "<span style='display:inline-block;width:8px;height:8px;border-radius:50%;background:#F87171;margin-right:5px;'></span>"
                    "<span style='display:inline-block;width:8px;height:8px;border-radius:50%;background:#FBBF24;margin-right:5px;'></span>"
                    "<span style='display:inline-block;width:8px;height:8px;border-radius:50%;background:#34D399;'></span>"
                    "</div>"
                    f"<pre style='background:#1E293B;color:#E2E8F0;padding:14px 16px;margin:0;"
                    f"overflow-x:auto;font-family:Consolas,Menlo,monospace;font-size:11.5px;line-height:1.5;"
                    f"white-space:pre-wrap;word-wrap:break-word;'><code>{code_text}</code></pre>"
                    "</div>"
                )
            continue
        if in_code_block:
            code_buffer.append(raw_line)
            continue

        # Skip the plain-text difficulty/audience line — it's rendered as pill badges instead
        if line_clean.startswith("**Difficulty:**"):
            continue

        # Escape HTML special chars in normal (non-code) content, then apply inline markdown
        escaped = html_lib.escape(line_clean)

        if line_clean.startswith("# "):
            close_list()
            html_lines.append(
                "<div style='background:linear-gradient(135deg,#1A2040,#2D3561);border-radius:12px;"
                "padding:32px 30px 24px 30px;margin-bottom:22px;border-bottom:5px solid #E9B259;'>"
                f"<h1 style='color:#FFFFFF;font-family:sans-serif;font-size:25px;margin:0;line-height:1.35;'>{_inline_markdown_to_html(html_lib.escape(line_clean[2:]))}</h1>"
                "<div style='color:#E9B259;font-family:sans-serif;font-size:11px;letter-spacing:1.5px;margin-top:10px;text-transform:uppercase;'>Curricula AI &middot; Course Material</div>"
                "</div>"
                "<div style='display:block;margin-bottom:20px;'>"
                f"<span style='display:inline-block;background:#FFF8EC;color:#C8913A;border:1px solid #E9B259;border-radius:20px;padding:4px 14px;font-family:sans-serif;font-size:10.5px;font-weight:bold;margin-right:6px;'>{html_lib.escape(str(_pdf_difficulty))}</span>"
                f"<span style='display:inline-block;background:#FFF8EC;color:#C8913A;border:1px solid #E9B259;border-radius:20px;padding:4px 14px;font-family:sans-serif;font-size:10.5px;font-weight:bold;margin-right:6px;'>{html_lib.escape(str(_pdf_audience))}</span>"
                f"<span style='display:inline-block;background:#2D3561;color:#FFFFFF;border-radius:20px;padding:4px 14px;font-family:sans-serif;font-size:10.5px;font-weight:bold;'>{html_lib.escape(str(role.upper()))} POV</span>"
                "</div>"
            )
        elif line_clean.startswith("## "):
            close_list()
            html_lines.append(
                "<div class='page-break' style='page-break-before:always;background:#FFF8EC;border-left:5px solid #E9B259;"
                "border-radius:6px;padding:12px 16px;margin-top:6px;'>"
                f"<h2 style='color:#2D3561;font-family:sans-serif;margin:0;font-size:17px;'>{_inline_markdown_to_html(html_lib.escape(line_clean[3:]))}</h2>"
                "</div>"
            )
        elif line_clean.startswith("### "):
            close_list()
            html_lines.append(
                "<h3 style='color:#FFFFFF;background:#C8913A;display:inline-block;font-family:sans-serif;"
                f"margin-top:22px;margin-bottom:6px;font-size:13.5px;padding:5px 12px;border-radius:5px;'>{_inline_markdown_to_html(html_lib.escape(line_clean[4:]))}</h3>"
            )
        elif line_clean.startswith("#### "):
            close_list()
            html_lines.append(
                "<h4 style='color:#2D3561;font-family:sans-serif;margin-top:14px;margin-bottom:4px;font-size:12.5px;"
                f"text-transform:uppercase;letter-spacing:0.5px;border-bottom:1.5px solid #E2E8F0;padding-bottom:4px;'>{_inline_markdown_to_html(html_lib.escape(line_clean[5:]))}</h4>"
            )
        elif raw_line.startswith("  - ") or raw_line.startswith("    - "):
            item = _inline_markdown_to_html(html_lib.escape(line_clean[2:]))
            if list_mode != "ul-nested":
                close_list()
                html_lines.append("<ul style='margin:4px 0 8px 0;padding-left:20px;list-style-type:circle;'><li style='list-style:none;'><ul style='margin:2px 0;padding-left:18px;list-style-type:circle;'>")
                list_mode = "ul-nested"
            html_lines.append(f"<li style='font-family:sans-serif;line-height:1.5;color:#475569;margin-bottom:3px;'>{item}</li>")
        elif line_clean.startswith("- "):
            item = _inline_markdown_to_html(html_lib.escape(line_clean[2:]))
            if list_mode != "ul":
                close_list()
                html_lines.append("<ul style='margin:4px 0 8px 0;padding-left:20px;'>")
                list_mode = "ul"
            html_lines.append(f"<li style='font-family:sans-serif;line-height:1.6;color:#334155;margin-bottom:4px;'>{item}</li>")
        elif re.match(r'^\d+\.\s', line_clean):
            close_list()
            html_lines.append(f"<p style='font-family:sans-serif;line-height:1.65;color:#334155;margin-bottom:8px;font-weight:600;'>{_inline_markdown_to_html(escaped)}</p>")
        elif line_clean == "---":
            close_list()
            html_lines.append("<hr style='border: none; border-top: 1px solid #E2E8F0; margin: 24px 0;'>")
        elif line_clean == "":
            close_list()
            html_lines.append("<div style='height:8px;'></div>")
        else:
            close_list()
            html_lines.append(f"<p style='font-family:sans-serif;line-height:1.65;color:#334155;margin-bottom:10px;'>{_inline_markdown_to_html(escaped)}</p>")

    close_list()
    body = "\n".join(html_lines)
    title = course_data.get('title', 'Exported Course')
    return f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        @page {{
            size: A4;
            margin: 20mm;
        }}
        body {{
            padding: 30px;
            max-width: 850px;
            margin: auto;
            color: #1A2040;
            background-color: #FFFFFF;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }}
        h1, h2, h3, h4 {{
            page-break-after: avoid;
        }}
        li {{
            page-break-inside: avoid;
        }}
        pre, div[style*="0F172A"] {{
            page-break-inside: avoid;
        }}
        @media print {{
            body {{ padding: 0; }}
            .page-break {{ page-break-before: always; }}
        }}
    </style>
</head>
<body>
    {body}
    <div style="text-align:center;font-size:9.5px;color:#94A3B8;border-top:1px solid #E2E8F0;padding-top:10px;margin-top:26px;font-family:sans-serif;">
        Curricula AI &middot; Maxy Academy
    </div>
</body>
</html>"""

def export_to_docx(course_data: dict, role: str) -> io.BytesIO:
    output = io.BytesIO()
    if not docx:
        txt = export_to_markdown(course_data, role)
        output.write(txt.encode('utf-8'))
        output.seek(0)
        return output

    doc = docx.Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(11)
    
    doc.add_heading(course_data.get('title', 'Untitled Course'), level=0)
    doc.add_paragraph(f"Difficulty: {course_data.get('config', {}).get('difficulty', 'Beginner')} | Audience: {course_data.get('config', {}).get('target_audience', 'Student')}")
    
    roles_to_export = ["creator", "student", "educator"] if role == "all" else [role]
    lessons = course_data.get("lessons", []) or course_data.get("structure", []) or [{"title": course_data.get("title", "Course Module")}]

    for r in roles_to_export:
        doc.add_heading(get_role_label(r), level=1)
        for idx, lesson in enumerate(lessons):
            l_num = lesson.get('order') or (idx + 1)
            clean_t = clean_lesson_title(lesson.get('title', 'Untitled Lesson'))
            doc.add_heading(f"Lesson {l_num}: {clean_t}", level=2)
            sections = get_resolved_lesson_sections(lesson, r)
            
            for sec_type, content in sections.items():
                doc.add_heading(sec_type.replace("_", " ").capitalize(), level=3)
                formatted_lines = format_section_content_to_md(content)
                for line in formatted_lines:
                    line_clean = line.strip()
                    if line_clean.startswith('- '):
                        doc.add_paragraph(line_clean[2:], style='List Bullet')
                    else:
                        doc.add_paragraph(line_clean)
                        
    doc.save(output)
    output.seek(0)
    return output

def md_to_reportlab_html(text: str) -> str:
    if not text:
        return ""
    s = html_lib.escape(text)
    s = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', s)
    s = re.sub(r'__(.+?)__', r'<b>\1</b>', s)
    s = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'<i>\1</i>', s)
    s = re.sub(r'(?<!_)_(?!_)(.+?)(?<!_)_(?!_)', r'<i>\1</i>', s)
    s = re.sub(r'`([^`]+?)`', r'<font name="Courier" color="#BE185D">\1</font>', s)
    return s

def _render_pdf_with_playwright(html_content: str) -> bytes:
    """Renders HTML to PDF bytes using headless Chromium. This supports the
    full modern CSS the new template needs (flexbox, grid, gradients) which
    wkhtmltopdf's old WebKit engine cannot render correctly."""
    with sync_playwright() as p:
        browser = p.chromium.launch()
        try:
            page = browser.new_page()
            page.set_content(html_content, wait_until="networkidle")
            pdf_bytes = page.pdf(
                format="A4",
                print_background=True,
                margin={"top": "0", "bottom": "0", "left": "0", "right": "0"},
            )
        finally:
            browser.close()
    return pdf_bytes


def export_to_pdf(course_data: dict, role: str) -> io.BytesIO:
    html_content = export_to_html(course_data, role)

    # 0. Preferred path: Playwright / headless Chromium (full CSS support —
    #    this is what makes the dark-theme cover/TOC/quiz-card/table layout
    #    render correctly; wkhtmltopdf below cannot).
    if sync_playwright:
        try:
            pdf_bytes = _render_pdf_with_playwright(html_content)
            out0 = io.BytesIO()
            out0.write(pdf_bytes)
            out0.seek(0)
            return out0
        except Exception as e:
            print(f"Playwright PDF render failed, falling back to wkhtmltopdf/reportlab: {e}")

    # Fallback chain below uses the legacy light-theme renderer, since
    # wkhtmltopdf/reportlab can't render the new template's CSS reliably.
    html_content = export_to_html_legacy(course_data, role)

    # 1. Find wkhtmltopdf binary path
    import os
    import subprocess
    import tempfile

    wkhtmltopdf_bin = "wkhtmltopdf"
    possible_paths = [
        r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe",
        r"C:\Program Files\wkhtmltopdf\wkhtmltopdf.exe",
        r"C:\Program Files (x86)\wkhtmltopdf\bin\wkhtmltopdf.exe",
        r"C:\Program Files (x86)\wkhtmltopdf\wkhtmltopdf.exe",
    ]
    for p in possible_paths:
        if os.path.exists(p):
            wkhtmltopdf_bin = p
            break

    wkhtmltopdf_options = {
        'quiet': '',
        'footer-center': 'Page [page] of [topage]',
        'footer-font-size': '8',
        'footer-font-name': 'Arial',
        'footer-text-color': '#94A3B8',
        'footer-spacing': '4',
    }

    # 1. Try pdfkit if available
    if pdfkit:
        try:
            config = None
            if wkhtmltopdf_bin != "wkhtmltopdf" and os.path.exists(wkhtmltopdf_bin):
                config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_bin)
            try:
                pdf_bytes = pdfkit.from_string(html_content, False, options=wkhtmltopdf_options, configuration=config)
            except Exception:
                # Footer options rejected by this wkhtmltopdf build — retry without them
                pdf_bytes = pdfkit.from_string(html_content, False, configuration=config)
            out1 = io.BytesIO()
            out1.write(pdf_bytes)
            out1.seek(0)
            return out1
        except Exception:
            pass

    # 2. Try direct wkhtmltopdf subprocess execution
    try:
        with tempfile.NamedTemporaryFile(suffix='.html', delete=False, mode='w', encoding='utf-8') as html_file:
            html_file.write(html_content)
            html_path = html_file.name
        
        pdf_path = html_path + '.pdf'
        footer_args = ['--footer-center', 'Page [page] of [topage]', '--footer-font-size', '8', '--footer-font-name', 'Arial']
        res = subprocess.run([wkhtmltopdf_bin] + footer_args + [html_path, pdf_path], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if res.returncode != 0 or not os.path.exists(pdf_path):
            # Retry without footer options in case this wkhtmltopdf build rejects them
            res = subprocess.run([wkhtmltopdf_bin, html_path, pdf_path], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if res.returncode == 0 and os.path.exists(pdf_path):
            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()
            out2 = io.BytesIO()
            out2.write(pdf_bytes)
            out2.seek(0)
            try:
                os.remove(html_path)
                os.remove(pdf_path)
            except Exception:
                pass
            return out2
    except Exception:
        pass

    # 3. ReportLab Multi-page Document Builder
    if SimpleDocTemplate:
        out3 = io.BytesIO()
        doc = SimpleDocTemplate(
            out3, 
            pagesize=letter,
            leftMargin=36,
            rightMargin=36,
            topMargin=40,
            bottomMargin=40
        )
        styles = getSampleStyleSheet()
        story = []
        
        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=18,
            leading=24,
            alignment=0,
            textColor=colors.HexColor('#FFFFFF'),
            backColor=colors.HexColor('#1A2040'),
            borderColor=colors.HexColor('#E9B259'),
            borderWidth=0,
            borderPadding=(16, 16, 16, 16),
            borderRadius=8,
            spaceAfter=14
        )
        h1_style = ParagraphStyle(
            'Heading1Custom',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=14,
            leading=17,
            alignment=0,
            textColor=colors.HexColor('#2D3561'),
            backColor=colors.HexColor('#FFF8EC'),
            borderPadding=(8, 8, 10, 10),
            borderRadius=5,
            spaceBefore=16,
            spaceAfter=10,
            keepWithNext=True
        )
        h2_style = ParagraphStyle(
            'Heading2Custom',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=11.5,
            leading=14,
            alignment=0,
            textColor=colors.HexColor('#FFFFFF'),
            backColor=colors.HexColor('#C8913A'),
            borderPadding=(5, 5, 8, 8),
            borderRadius=4,
            spaceBefore=12,
            spaceAfter=7,
            keepWithNext=True
        )
        h3_style = ParagraphStyle(
            'Heading3Custom',
            parent=styles['Heading3'],
            fontName='Helvetica-Bold',
            fontSize=10.5,
            leading=13,
            textColor=colors.HexColor('#334155'),
            spaceBefore=8,
            spaceAfter=4,
            keepWithNext=True
        )
        body_style = ParagraphStyle(
            'BodyCustom',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=9.5,
            leading=13.5,
            textColor=colors.HexColor('#334155'),
            spaceAfter=5
        )
        bullet_style_1 = ParagraphStyle(
            'Bullet1',
            parent=body_style,
            leftIndent=15,
            firstLineIndent=-10,
            spaceAfter=3
        )
        bullet_style_2 = ParagraphStyle(
            'Bullet2',
            parent=body_style,
            leftIndent=30,
            firstLineIndent=-10,
            spaceAfter=2
        )
        code_style = ParagraphStyle(
            'CodeBlockStyle',
            parent=body_style,
            fontName='Courier',
            fontSize=8.5,
            leading=11,
            textColor=colors.HexColor('#E2E8F0'),
            backColor=colors.HexColor('#1E293B'),
            borderPadding=8,
            spaceBefore=6,
            spaceAfter=6
        )

        txt = export_to_markdown(course_data, role)
        
        # Clean emojis for ReportLab standard Helvetica font
        txt = re.sub(r'[\U00010000-\U0010ffff]', '', txt)
        txt = re.sub(r'[\u2600-\u27BF]', '', txt)

        in_code = False
        code_lines = []

        for raw_line in txt.split('\n'):
            line_str = raw_line.rstrip()
            line_clean = line_str.strip()

            if line_clean.startswith('```'):
                if not in_code:
                    in_code = True
                    code_lines = []
                else:
                    in_code = False
                    code_text = html_lib.escape("\n".join(code_lines))
                    story.append(Paragraph(code_text.replace("\n", "<br/>").replace(" ", "&nbsp;"), code_style))
                continue

            if in_code:
                code_lines.append(raw_line)
                continue

            if not line_clean:
                story.append(Spacer(1, 4))
                continue

            if line_clean.startswith('# '):
                story.append(Paragraph(md_to_reportlab_html(line_clean[2:]), title_style))
            elif line_clean.startswith('## '):
                story.append(Paragraph(md_to_reportlab_html(line_clean[3:]), h1_style))
            elif line_clean.startswith('### '):
                story.append(Paragraph(md_to_reportlab_html(line_clean[4:]), h2_style))
            elif line_clean.startswith('#### '):
                story.append(Paragraph(md_to_reportlab_html(line_clean[5:]), h3_style))
            elif line_clean == '---':
                try:
                    from reportlab.platypus import HRFlowable
                    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#CBD5E1'), spaceBefore=8, spaceAfter=8))
                except Exception:
                    story.append(Spacer(1, 8))
            elif line_str.startswith('    - ') or line_str.startswith('\t- '):
                item_text = line_clean[2:].strip()
                story.append(Paragraph(f"&bull;&nbsp;{md_to_reportlab_html(item_text)}", bullet_style_2))
            elif line_str.startswith('  - '):
                item_text = line_clean[2:].strip()
                story.append(Paragraph(f"&bull;&nbsp;{md_to_reportlab_html(item_text)}", bullet_style_1))
            elif line_clean.startswith('- ') or line_clean.startswith('* '):
                item_text = line_clean[2:].strip()
                story.append(Paragraph(f"&bull;&nbsp;{md_to_reportlab_html(item_text)}", bullet_style_1))
            else:
                story.append(Paragraph(md_to_reportlab_html(line_clean), body_style))

        def _draw_page_decor(canvas_obj, doc_obj):
            try:
                canvas_obj.saveState()
                width, height = letter
                # Thin accent bar across the top of every page
                canvas_obj.setFillColor(colors.HexColor('#E9B259'))
                canvas_obj.rect(0, height - 6, width, 6, fill=1, stroke=0)
                # Footer: page number + brand
                canvas_obj.setFont('Helvetica', 8)
                canvas_obj.setFillColor(colors.HexColor('#94A3B8'))
                canvas_obj.drawCentredString(width / 2.0, 24, f"Page {doc_obj.page}")
                canvas_obj.drawString(36, 24, "Curricula AI")
                canvas_obj.drawRightString(width - 36, 24, "Maxy Academy")
                canvas_obj.restoreState()
            except Exception:
                pass

        try:
            doc.build(story, onFirstPage=_draw_page_decor, onLaterPages=_draw_page_decor)
            out3.seek(0)
            return out3
        except Exception as e:
            print(f"ReportLab error: {e}")

    # Fallback minimal Canvas document
    try:
        from reportlab.pdfgen import canvas
        out4 = io.BytesIO()
        c = canvas.Canvas(out4, pagesize=letter)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, 750, course_data.get('title', 'Exported Course'))
        c.setFont("Helvetica", 10)
        c.drawString(50, 730, f"Role: {role.upper()} POV | Document Ready")
        c.save()
        out4.seek(0)
        return out4
    except Exception:
        pass

    out5 = io.BytesIO()
    pdf_skeleton = b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R>>endobj\nxref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000052 00000 n \n0000000101 00000 n \ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n178\n%%EOF"
    out5.write(pdf_skeleton)
    out5.seek(0)
    return out5




def export_all_zip(course_data: dict) -> io.BytesIO:
    output = io.BytesIO()
    with zipfile.ZipFile(output, 'w') as zip_file:
        for role in ["creator", "student", "educator"]:
            # PDF Document
            try:
                pdf_stream = export_to_pdf(course_data, role)
                zip_file.writestr(f"{role}_pov.pdf", pdf_stream.getvalue())
            except Exception as e:
                print(f"Error generating ZIP PDF for {role}: {e}")

            # DOCX Document
            try:
                docx_stream = export_to_docx(course_data, role)
                zip_file.writestr(f"{role}_pov.docx", docx_stream.getvalue())
            except Exception as e:
                print(f"Error generating ZIP DOCX for {role}: {e}")

            # HTML Web Page Document
            try:
                html_content = export_to_html(course_data, role)
                zip_file.writestr(f"{role}_pov.html", html_content)
            except Exception as e:
                print(f"Error generating ZIP HTML for {role}: {e}")

            # Markdown Document
            try:
                md_content = export_to_markdown(course_data, role)
                zip_file.writestr(f"{role}_pov.md", md_content)
            except Exception as e:
                print(f"Error generating ZIP MD for {role}: {e}")
            
    output.seek(0)
    return output


def _hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_pptx_from_structure(slides_json: dict, layout: str = "modern", brand_colors: dict = None) -> io.BytesIO:
    """Create a PPTX file from AI-generated slide structure."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        output = io.BytesIO()
        output.write(b"python-pptx not installed")
        output.seek(0)
        return output

    layout_data = slides_json.get("layouts", {}).get(layout, slides_json.get("layouts", {}).get("modern", {}))
    theme = layout_data.get("theme", {})
    slides_data = layout_data.get("slides", [])

    if brand_colors:
        theme.update(brand_colors)

    primary_rgb = RGBColor(*_hex_to_rgb(theme.get("primary", "#1a202c")))
    secondary_rgb = RGBColor(*_hex_to_rgb(theme.get("secondary", "#ffffff")))
    accent_rgb = RGBColor(*_hex_to_rgb(theme.get("accent", "#d69e2e")))
    text_rgb = RGBColor(*_hex_to_rgb(theme.get("text", "#1a202c")))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = primary_rgb

        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data.get("notes", "")

        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            _add_title_slide(slide, slide_data, primary_rgb, secondary_rgb, accent_rgb, layout)
        elif slide_type == "toc":
            _add_toc_slide(slide, slide_data, primary_rgb, secondary_rgb, accent_rgb, text_rgb, layout)
        elif slide_type == "lesson_title":
            _add_lesson_title_slide(slide, slide_data, primary_rgb, secondary_rgb, accent_rgb, layout)
        elif slide_type == "content":
            _add_content_slide(slide, slide_data, primary_rgb, secondary_rgb, accent_rgb, text_rgb, layout)
        elif slide_type == "code":
            _add_code_slide(slide, slide_data, primary_rgb, secondary_rgb, accent_rgb, text_rgb, layout)
        elif slide_type == "end":
            _add_end_slide(slide, slide_data, primary_rgb, secondary_rgb, accent_rgb, layout)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _add_title_slide(slide, data, primary_rgb, secondary_rgb, accent_rgb, layout="layout_1"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = primary_rgb

    if layout == "layout_2":
        shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape1.line.fill.background()
        circle = slide.shapes.add_shape(9, Inches(9), Inches(-1), Inches(5.5), Inches(5.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb
        circle.fill.fore_color.brightness = 0.7
        circle.line.fill.background()
    elif layout == "layout_3":
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Course Title")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = secondary_rgb
    p.alignment = PP_ALIGN.CENTER

    subtitle = data.get("subtitle", "")
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf2 = subtitle_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = accent_rgb
        p2.alignment = PP_ALIGN.CENTER

    if layout == "layout_1":
        accent_line = slide.shapes.add_shape(1, Inches(4), Inches(4.5), Inches(5), Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_2":
        accent_line = slide.shapes.add_shape(1, Inches(5), Inches(4.5), Inches(3.333), Inches(0.04))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_3":
        accent_line = slide.shapes.add_shape(1, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.03))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()


def _add_toc_slide(slide, data, primary_rgb, secondary_rgb, accent_rgb, text_rgb, layout="layout_1"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if layout == "layout_2":
        shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape1.line.fill.background()
        circle = slide.shapes.add_shape(9, Inches(-2), Inches(4), Inches(6), Inches(6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb
        circle.fill.fore_color.brightness = 0.75
        circle.line.fill.background()
    elif layout == "layout_3":
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "Table of Contents")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = secondary_rgb

    if layout == "layout_1":
        accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_2":
        accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(4), Inches(0.03))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_3":
        accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(2), Inches(0.02))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()

    items = data.get("items", [])
    content_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = text_rgb
        p.space_after = Pt(10)
        if layout == "layout_1":
            mark = slide.shapes.add_shape(1, Inches(0.8), Inches(1.9 + i * 0.55), Inches(0.18), Inches(0.18))
            mark.fill.solid()
            mark.fill.fore_color.rgb = accent_rgb
            mark.line.fill.background()
        elif layout == "layout_2":
            mark = slide.shapes.add_shape(9, Inches(0.85), Inches(1.95 + i * 0.55), Inches(0.14), Inches(0.14))
            mark.fill.solid()
            mark.fill.fore_color.rgb = accent_rgb
            mark.line.fill.background()
        elif layout == "layout_3":
            pass


def _add_lesson_title_slide(slide, data, primary_rgb, secondary_rgb, accent_rgb, layout="layout_1"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if layout == "layout_1":
        accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(0.3), Inches(2.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_rgb
        accent_bar.line.fill.background()
    elif layout == "layout_2":
        shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape1.line.fill.background()
        circle = slide.shapes.add_shape(9, Inches(8.5), Inches(1), Inches(6), Inches(6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb
        circle.fill.fore_color.brightness = 0.7
        circle.line.fill.background()
    elif layout == "layout_3":
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1) if layout != "layout_1" else Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Lesson")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = secondary_rgb

    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = accent_rgb

    if layout == "layout_3":
        bottom_line = slide.shapes.add_shape(1, Inches(1), Inches(4), Inches(2), Inches(0.02))
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = accent_rgb
        bottom_line.line.fill.background()


def _add_content_slide(slide, data, primary_rgb, secondary_rgb, accent_rgb, text_rgb, layout="layout_1"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if layout == "layout_2":
        shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape1.line.fill.background()
        circle = slide.shapes.add_shape(9, Inches(10), Inches(-2), Inches(5), Inches(5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb
        circle.fill.fore_color.brightness = 0.75
        circle.line.fill.background()
    elif layout == "layout_3":
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "Content")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = secondary_rgb

    if layout == "layout_1":
        accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.04))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_2":
        accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.03))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_3":
        accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.02))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()

    bullets = data.get("bullets", [])
    content_box = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(11), Inches(5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = text_rgb
        p.space_after = Pt(8)

        if layout == "layout_1":
            mark = slide.shapes.add_shape(1, Inches(1.05), Inches(1.9 + i * 0.6), Inches(0.15), Inches(0.15))
            mark.fill.solid()
            mark.fill.fore_color.rgb = accent_rgb
            mark.line.fill.background()
        elif layout == "layout_2":
            mark = slide.shapes.add_shape(9, Inches(1.05), Inches(1.95 + i * 0.6), Inches(0.12), Inches(0.12))
            mark.fill.solid()
            mark.fill.fore_color.rgb = accent_rgb
            mark.line.fill.background()
        elif layout == "layout_3":
            pass


def _add_code_slide(slide, data, primary_rgb, secondary_rgb, accent_rgb, text_rgb, layout="layout_1"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if layout == "layout_2":
        shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape1.line.fill.background()
    elif layout == "layout_3":
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "Code Example")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = secondary_rgb

    if layout == "layout_1":
        code_bg = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = RGBColor(30, 30, 40)
        code_bg.line.color.rgb = RGBColor(60, 60, 80)
    elif layout == "layout_2":
        code_bg = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = RGBColor(15, 25, 45)
        code_bg.line.color.rgb = accent_rgb
    elif layout == "layout_3":
        code_bg = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = RGBColor(240, 240, 245)
        code_bg.line.color.rgb = RGBColor(200, 200, 210)

    code_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.8), Inches(4.6))
    tf2 = code_box.text_frame
    tf2.word_wrap = True
    code_text = data.get("code", "# Code here")
    for i, line in enumerate(code_text.split('\n')):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Courier New"
        if layout == "layout_3":
            p.font.color.rgb = RGBColor(40, 40, 60)
        else:
            p.font.color.rgb = RGBColor(0, 200, 120)


def _add_end_slide(slide, data, primary_rgb, secondary_rgb, accent_rgb, layout="layout_1"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if layout == "layout_2":
        shape1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape1.line.fill.background()
        circle = slide.shapes.add_shape(9, Inches(4), Inches(1), Inches(5.333), Inches(5.333))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb
        circle.fill.fore_color.brightness = 0.7
        circle.line.fill.background()
    elif layout == "layout_3":
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "Thank You")
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = secondary_rgb
    p.alignment = PP_ALIGN.CENTER

    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = accent_rgb
        p2.alignment = PP_ALIGN.CENTER

    if layout == "layout_1":
        accent_line = slide.shapes.add_shape(1, Inches(4), Inches(3.9), Inches(5), Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_2":
        accent_line = slide.shapes.add_shape(1, Inches(5), Inches(3.9), Inches(3.333), Inches(0.04))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()
    elif layout == "layout_3":
        accent_line = slide.shapes.add_shape(1, Inches(5.5), Inches(3.95), Inches(2.333), Inches(0.03))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_rgb
        accent_line.line.fill.background()